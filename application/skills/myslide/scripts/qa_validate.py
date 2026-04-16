#!/usr/bin/env python3
"""
PPTX Layout & Design QA Validator

Catches structural issues that visual inspection misses because rendered
images clip objects at slide boundaries.

Problem: When a shape extends beyond the slide edge, soffice/pdftoppm
simply crops it. The QA subagent sees a normal-looking image and passes it.
This script reads the OOXML directly and flags what the renderer hides.

Checks:
  - bounds:     Shapes/groups outside slide boundaries
  - connector:  Arrow/line endpoints outside slide, zero-length connectors
  - font_size:  Text below 15pt minimum (body) or 8pt (caption)
  - zero_size:  Invisible shapes (0 width or 0 height)

Usage:
    python3 qa_validate.py presentation.pptx
    python3 qa_validate.py presentation.pptx --json
    python3 qa_validate.py presentation.pptx --strict

Exit codes:
    0  No critical issues found
    1  Critical issues found — must fix before delivery
    2  Error (file not found, parse failure, etc.)
"""

import argparse
import json
import sys
from collections import defaultdict
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    import subprocess

    subprocess.check_call(
        [sys.executable, "-m", "pip", "install", "python-pptx", "lxml", "-q"]
    )
    from pptx import Presentation


# ── Constants ──────────────────────────────────────────────────────────

EMU_PER_INCH = 914400

# Font thresholds in 100ths-of-a-point (OOXML stores sizes this way).
# 1500 = 15pt (body minimum), 800 = 8pt (caption minimum).
MIN_BODY_SZ = 1500
MIN_CAPTION_SZ = 800

# Shapes covering >90% of slide area are treated as backgrounds/overlays
# and excluded from bounds checking (they're intentionally full-bleed).
BG_AREA_RATIO = 0.90

# OOXML namespaces
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ── Issue ──────────────────────────────────────────────────────────────


class Issue:
    def __init__(self, level, slide_num, check, message, **details):
        self.level = level  # "critical" | "warning" | "info"
        self.slide_num = slide_num
        self.check = check
        self.message = message
        self.details = details

    def to_dict(self):
        d = {
            "level": self.level,
            "slide": self.slide_num,
            "check": self.check,
            "message": self.message,
        }
        d.update(self.details)
        return d

    def __str__(self):
        tag = self.level.upper()
        return f"[{tag}] {self.message}"


# ── Helpers ────────────────────────────────────────────────────────────


def emu_to_in(emu):
    """EMU → inches, 2 decimal places."""
    return round(emu / EMU_PER_INCH, 2) if emu is not None else None


def _is_bg(shape, sw, sh):
    """True if shape covers >90% of slide (background/overlay)."""
    w, h = shape.width, shape.height
    if w is None or h is None:
        return False
    return (w * h) > (sw * sh * BG_AREA_RATIO)


def _font_sz(run_elem):
    """Return font size in 100ths-pt from <a:rPr sz='...'>, or None."""
    rPr = run_elem.find(f"{{{A_NS}}}rPr")
    if rPr is not None:
        sz = rPr.get("sz")
        if sz is not None:
            return int(sz)
    return None


def _shape_name_from_cxn(cxn_elem):
    """Extract shape name from a <p:cxnSp> element."""
    nvPr = cxn_elem.find(f"{{{P_NS}}}nvCxnSpPr")
    if nvPr is not None:
        cNvPr = nvPr.find(f"{{{P_NS}}}cNvPr")
        if cNvPr is not None:
            return cNvPr.get("name", "Connector")
    return "Connector"


# ── Check: bounds ──────────────────────────────────────────────────────


def check_bounds(prs):
    """Flag shapes whose bounding box extends beyond the slide."""
    issues = []
    sw, sh = prs.slide_width, prs.slide_height

    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            if _is_bg(shape, sw, sh):
                continue

            l, t = shape.left, shape.top
            w = shape.width or 0
            h = shape.height or 0
            r, b = l + w, t + h

            overflows = []
            if l < 0:
                overflows.append(f"left by {emu_to_in(abs(l))}\"")
            if t < 0:
                overflows.append(f"top by {emu_to_in(abs(t))}\"")
            if r > sw:
                overflows.append(f"right by {emu_to_in(r - sw)}\"")
            if b > sh:
                overflows.append(f"bottom by {emu_to_in(b - sh)}\"")

            if not overflows:
                continue

            fully_outside = r <= 0 or b <= 0 or l >= sw or t >= sh
            verb = "entirely outside" if fully_outside else "extends beyond"

            issues.append(
                Issue(
                    "critical",
                    idx,
                    "bounds",
                    f"Shape '{shape.name}' {verb} slide ({', '.join(overflows)})",
                    shape=shape.name,
                    position=f"({emu_to_in(l)}\", {emu_to_in(t)}\")",
                    size=f"({emu_to_in(w)}\" x {emu_to_in(h)}\")",
                )
            )

    return issues


# ── Check: connectors ─────────────────────────────────────────────────


def check_connectors(prs):
    """Validate connector/arrow endpoints and lengths."""
    issues = []
    sw, sh = prs.slide_width, prs.slide_height

    for idx, slide in enumerate(prs.slides, 1):
        slide_xml = slide._element

        for cxn in slide_xml.findall(f".//{{{P_NS}}}cxnSp"):
            name = _shape_name_from_cxn(cxn)

            xfrm = cxn.find(f".//{{{A_NS}}}xfrm")
            if xfrm is None:
                continue

            off = xfrm.find(f"{{{A_NS}}}off")
            ext = xfrm.find(f"{{{A_NS}}}ext")
            if off is None or ext is None:
                continue

            x = int(off.get("x", 0))
            y = int(off.get("y", 0))
            cx = int(ext.get("cx", 0))
            cy = int(ext.get("cy", 0))

            # Zero-length connector → invisible
            if cx == 0 and cy == 0:
                issues.append(
                    Issue(
                        "warning",
                        idx,
                        "connector",
                        f"Connector '{name}' has zero length (invisible)",
                        shape=name,
                    )
                )
                continue

            # Determine actual start/end considering flips
            flip_h = xfrm.get("flipH") == "1"
            flip_v = xfrm.get("flipV") == "1"

            sx = (x + cx) if flip_h else x
            sy = (y + cy) if flip_v else y
            ex = x if flip_h else (x + cx)
            ey = y if flip_v else (y + cy)

            for label, px, py in [("start", sx, sy), ("end", ex, ey)]:
                if px < 0 or py < 0 or px > sw or py > sh:
                    issues.append(
                        Issue(
                            "critical",
                            idx,
                            "connector",
                            f"Connector '{name}' {label} point outside slide "
                            f"at ({emu_to_in(px)}\", {emu_to_in(py)}\")",
                            shape=name,
                        )
                    )

            # Very short connector (< 0.1") — often a positioning mistake
            length_sq = cx * cx + cy * cy
            min_len = int(0.1 * EMU_PER_INCH)
            if 0 < length_sq < min_len * min_len:
                length_in = emu_to_in(int(length_sq**0.5))
                issues.append(
                    Issue(
                        "info",
                        idx,
                        "connector",
                        f"Connector '{name}' is very short ({length_in}\")",
                        shape=name,
                    )
                )

    return issues


# ── Check: font sizes ─────────────────────────────────────────────────


def check_font_sizes(prs):
    """Flag text below the 15pt body minimum or 8pt caption minimum."""
    issues = []

    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    sz = _font_sz(run._r)
                    if sz is None:
                        continue  # inherited from theme — can't check

                    text = run.text.strip()
                    if not text:
                        continue

                    preview = text[:40]
                    pt = sz / 100

                    if sz < MIN_CAPTION_SZ:
                        issues.append(
                            Issue(
                                "critical",
                                idx,
                                "font_size",
                                f"Text {pt}pt in '{shape.name}' (min 8pt): \"{preview}\"",
                                shape=shape.name,
                                size_pt=pt,
                            )
                        )
                    elif sz < MIN_BODY_SZ:
                        issues.append(
                            Issue(
                                "warning",
                                idx,
                                "font_size",
                                f"Text {pt}pt in '{shape.name}' (body min 15pt): \"{preview}\"",
                                shape=shape.name,
                                size_pt=pt,
                            )
                        )

    return issues


# ── Check: zero-size shapes ───────────────────────────────────────────


def check_zero_size(prs):
    """Flag shapes with zero width or height (invisible/broken)."""
    issues = []

    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            w, h = shape.width, shape.height
            if w is not None and h is not None and (w == 0 or h == 0):
                dim = "width" if w == 0 else "height"
                issues.append(
                    Issue(
                        "warning",
                        idx,
                        "zero_size",
                        f"Shape '{shape.name}' has zero {dim} (invisible)",
                        shape=shape.name,
                    )
                )

    return issues


# ── Orchestration ──────────────────────────────────────────────────────


def validate(pptx_path, strict=False):
    """Run all checks. Returns (issues, presentation)."""
    prs = Presentation(pptx_path)

    all_issues = []
    all_issues.extend(check_bounds(prs))
    all_issues.extend(check_connectors(prs))
    all_issues.extend(check_font_sizes(prs))
    all_issues.extend(check_zero_size(prs))

    if not strict:
        all_issues = [i for i in all_issues if i.level != "info"]

    return all_issues, prs


def format_report(issues, prs, pptx_path):
    """Human-readable report."""
    sw = emu_to_in(prs.slide_width)
    sh = emu_to_in(prs.slide_height)
    n = len(prs.slides)

    lines = [
        "=== PPTX QA Validation Report ===",
        f"File: {Path(pptx_path).name}",
        f"Slides: {n}   Size: {sw}\" x {sh}\"",
        "",
    ]

    if not issues:
        lines.append("All checks passed — no issues found.")
        return "\n".join(lines)

    by_slide = defaultdict(list)
    for issue in issues:
        by_slide[issue.slide_num].append(issue)

    for sn in sorted(by_slide):
        lines.append(f"--- Slide {sn} ---")
        for issue in by_slide[sn]:
            lines.append(f"  {issue}")
            # Add position/size details for bounds issues
            if "position" in issue.details:
                lines.append(
                    f"    Position: {issue.details['position']}  "
                    f"Size: {issue.details['size']}"
                )
        lines.append("")

    crit = sum(1 for i in issues if i.level == "critical")
    warn = sum(1 for i in issues if i.level == "warning")
    info = sum(1 for i in issues if i.level == "info")

    lines.append("=== Summary ===")
    parts = []
    if crit:
        parts.append(f"{crit} critical")
    if warn:
        parts.append(f"{warn} warning")
    if info:
        parts.append(f"{info} info")
    lines.append(f"Total: {len(issues)} ({', '.join(parts)})")
    lines.append(f"Slides with issues: {len(by_slide)} / {n}")

    return "\n".join(lines)


# ── CLI ────────────────────────────────────────────────────────────────


def main():
    parser = argparse.ArgumentParser(
        description="PPTX Layout & Design QA Validator"
    )
    parser.add_argument("pptx", help="Path to .pptx file")
    parser.add_argument("--json", action="store_true", help="Output JSON")
    parser.add_argument(
        "--strict", action="store_true", help="Include INFO-level findings"
    )
    args = parser.parse_args()

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"Error: {pptx_path} not found", file=sys.stderr)
        sys.exit(2)

    try:
        issues, prs = validate(str(pptx_path), strict=args.strict)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(2)

    if args.json:
        out = {
            "file": str(pptx_path),
            "slides": len(prs.slides),
            "slide_width_in": emu_to_in(prs.slide_width),
            "slide_height_in": emu_to_in(prs.slide_height),
            "issues": [i.to_dict() for i in issues],
            "summary": {
                "total": len(issues),
                "critical": sum(1 for i in issues if i.level == "critical"),
                "warning": sum(1 for i in issues if i.level == "warning"),
                "info": sum(1 for i in issues if i.level == "info"),
            },
        }
        print(json.dumps(out, indent=2, ensure_ascii=False))
    else:
        print(format_report(issues, prs, pptx_path))

    has_critical = any(i.level == "critical" for i in issues)
    sys.exit(1 if has_critical else 0)


if __name__ == "__main__":
    main()
