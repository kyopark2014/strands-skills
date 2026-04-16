#!/usr/bin/env python3
"""
PPTX Animation Engine — inject OOXML animations from a JSON spec.

Uses the correct PowerPoint OOXML structure where:
- presetID/presetClass/presetSubtype go on the p:cTn wrapper
- p:animEffect uses transition/filter attributes with p:cBhvr targeting
- Entrance animations include a p:set for style.visibility
- Targets are specified via p:cBhvr > p:tgtEl > p:spTgt

Usage:
    python3 apply_animations.py input.pptx animations.json -o output.pptx
    python3 apply_animations.py input.pptx --list-shapes
    python3 apply_animations.py input.pptx --inspect
"""

import argparse
import json
import sys
import subprocess

try:
    from pptx import Presentation
    from lxml import etree
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "lxml", "-q"])
    from pptx import Presentation
    from lxml import etree

# OOXML namespaces
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

NSMAP = {"p": P_NS, "a": A_NS, "r": R_NS}


def _qn(tag):
    """Convert 'p:timing' → '{ns}timing'."""
    prefix, local = tag.split(":")
    return f"{{{NSMAP[prefix]}}}{local}"


class IdCounter:
    def __init__(self, start=1):
        self.val = start

    def next(self):
        v = self.val
        self.val += 1
        return str(v)


# ---------------------------------------------------------------------------
# Effect Registry — maps effect keys to OOXML presetID + filter strings
# ---------------------------------------------------------------------------

# Filter strings used by p:animEffect
EFFECT_FILTERS = {
    # Entrance
    "appear":    {"presetID": "1",  "presetClass": "entr", "filter": None},
    "fade_in":   {"presetID": "10", "presetClass": "entr", "filter": "fade"},
    "fly_in":    {"presetID": "2",  "presetClass": "entr", "filter": "fly"},
    "wipe":      {"presetID": "22", "presetClass": "entr", "filter": "wipe({dir})"},
    "zoom_in":   {"presetID": "53", "presetClass": "entr", "filter": "wheel(1)"},
    "float_in":  {"presetID": "42", "presetClass": "entr", "filter": "fade"},
    # Exit
    "disappear": {"presetID": "1",  "presetClass": "exit", "filter": None},
    "fade_out":  {"presetID": "10", "presetClass": "exit", "filter": "fade"},
    "fly_out":   {"presetID": "2",  "presetClass": "exit", "filter": "fly"},
    # Emphasis
    "pulse":     {"presetID": "26", "presetClass": "emph", "filter": None},
    # Motion
    "motion_path": {"presetID": "0", "presetClass": "path", "filter": None},
}

DIRECTION_SUBTYPE = {
    "left": "4", "right": "2", "top": "1", "bottom": "8",
    "top_left": "5", "top_right": "3", "bottom_left": "9", "bottom_right": "6",
}

WIPE_DIR = {"left": "left", "right": "right", "top": "up", "bottom": "down"}


# ---------------------------------------------------------------------------
# Core: Build a single animation's children (p:set + p:animEffect or p:animMotion)
# ---------------------------------------------------------------------------

def _build_visibility_set(ids, sp_id):
    """Build p:set that changes style.visibility to 'visible'."""
    p_set = etree.Element(_qn("p:set"))

    ctn = etree.SubElement(p_set, _qn("p:cTn"), attrib={
        "id": ids.next(), "dur": "1", "fill": "hold"
    })
    st = etree.SubElement(ctn, _qn("p:stCondLst"))
    etree.SubElement(st, _qn("p:cond"), attrib={"delay": "0"})

    cbhvr = etree.SubElement(p_set, _qn("p:cBhvr"))
    cbhvr_ctn = etree.SubElement(cbhvr, _qn("p:cTn"), attrib={
        "id": ids.next(), "dur": "1", "fill": "hold"
    })
    cbhvr_st = etree.SubElement(cbhvr_ctn, _qn("p:stCondLst"))
    etree.SubElement(cbhvr_st, _qn("p:cond"), attrib={"delay": "0"})

    tgt_el = etree.SubElement(cbhvr, _qn("p:tgtEl"))
    etree.SubElement(tgt_el, _qn("p:spTgt"), attrib={"spid": str(sp_id)})

    attr_lst = etree.SubElement(cbhvr, _qn("p:attrNameLst"))
    attr_name = etree.SubElement(attr_lst, _qn("p:attrName"))
    attr_name.text = "style.visibility"

    to_elem = etree.SubElement(p_set, _qn("p:to"))
    str_val = etree.SubElement(to_elem, _qn("p:strVal"), attrib={"val": "visible"})

    return p_set


def _build_anim_effect(ids, sp_id, duration, transition, filter_str):
    """Build p:animEffect with p:cBhvr targeting."""
    anim_effect = etree.Element(_qn("p:animEffect"), attrib={
        "transition": transition, "filter": filter_str
    })

    cbhvr = etree.SubElement(anim_effect, _qn("p:cBhvr"))
    ctn = etree.SubElement(cbhvr, _qn("p:cTn"), attrib={
        "id": ids.next(), "dur": str(duration)
    })
    tgt_el = etree.SubElement(cbhvr, _qn("p:tgtEl"))
    etree.SubElement(tgt_el, _qn("p:spTgt"), attrib={"spid": str(sp_id)})

    return anim_effect


def _build_anim_motion(ids, sp_id, duration, path):
    """Build p:animMotion with p:cBhvr targeting."""
    anim_motion = etree.Element(_qn("p:animMotion"), attrib={
        "origin": "layout", "path": path, "pathEditMode": "relative"
    })

    cbhvr = etree.SubElement(anim_motion, _qn("p:cBhvr"))
    ctn = etree.SubElement(cbhvr, _qn("p:cTn"), attrib={
        "id": ids.next(), "dur": str(duration), "fill": "hold"
    })
    tgt_el = etree.SubElement(cbhvr, _qn("p:tgtEl"))
    etree.SubElement(tgt_el, _qn("p:spTgt"), attrib={"spid": str(sp_id)})

    return anim_motion


def _build_anim_scale(ids, sp_id, duration, from_pct=0, to_pct=100):
    """Build p:animScale for zoom effects."""
    anim_scale = etree.Element(_qn("p:animScale"))

    cbhvr = etree.SubElement(anim_scale, _qn("p:cBhvr"))
    ctn = etree.SubElement(cbhvr, _qn("p:cTn"), attrib={
        "id": ids.next(), "dur": str(duration), "fill": "hold"
    })
    tgt_el = etree.SubElement(cbhvr, _qn("p:tgtEl"))
    etree.SubElement(tgt_el, _qn("p:spTgt"), attrib={"spid": str(sp_id)})

    etree.SubElement(anim_scale, _qn("p:by"), attrib={
        "x": str(to_pct * 1000), "y": str(to_pct * 1000)
    })

    return anim_scale


# ---------------------------------------------------------------------------
# Build the innermost p:par that wraps one animation effect
# (the p:cTn here carries presetID, presetClass, etc.)
# ---------------------------------------------------------------------------

def build_effect_par(ids, sp_id, anim_spec, grp_id=0):
    """
    Build one animation as a p:par containing:
      p:cTn (presetID, presetClass, grpId, nodeType)
        p:stCondLst > p:cond
        p:childTnLst >
          p:set (visibility) [for entrance]
          p:animEffect / p:animMotion / p:animScale
    """
    effect = anim_spec["effect"]
    duration = anim_spec.get("duration", 500)
    direction = anim_spec.get("direction", "left")

    info = EFFECT_FILTERS.get(effect)
    if not info:
        print(f"WARNING: Unknown effect '{effect}'")
        return None

    preset_class = info["presetClass"]
    preset_id = info["presetID"]

    # Resolve subtype from direction
    if effect in ("fly_in", "fly_out"):
        subtype = DIRECTION_SUBTYPE.get(direction, "4")
    elif effect == "wipe":
        subtype = DIRECTION_SUBTYPE.get(direction, "4")
    else:
        subtype = "0"

    # Determine fill
    fill = "hold" if preset_class != "exit" else "hold"

    # Outer p:par
    par = etree.Element(_qn("p:par"))
    ctn = etree.SubElement(par, _qn("p:cTn"), attrib={
        "id": ids.next(),
        "presetID": preset_id,
        "presetClass": preset_class,
        "presetSubtype": subtype,
        "fill": fill,
        "grpId": str(grp_id),
        "nodeType": "clickEffect",
    })
    st = etree.SubElement(ctn, _qn("p:stCondLst"))
    etree.SubElement(st, _qn("p:cond"), attrib={"delay": "0"})
    child_tn = etree.SubElement(ctn, _qn("p:childTnLst"))

    # For entrance effects: add visibility set
    if preset_class == "entr":
        vis_set = _build_visibility_set(ids, sp_id)
        child_tn.append(vis_set)

    # Add the actual effect element
    if effect == "motion_path":
        path = anim_spec.get("path", "M 0 0 L 0.25 0 E")
        motion = _build_anim_motion(ids, sp_id, duration, path)
        child_tn.append(motion)

    elif effect == "zoom_in":
        # Zoom uses animScale
        scale = _build_anim_scale(ids, sp_id, duration, from_pct=0, to_pct=100)
        child_tn.append(scale)
        # Also add visibility
        if preset_class != "entr":
            vis_set = _build_visibility_set(ids, sp_id)
            child_tn.insert(0, vis_set)

    elif effect in ("appear", "disappear"):
        # Appear/Disappear only needs the visibility set (already added above for entr)
        if preset_class == "exit":
            vis_set = _build_visibility_set(ids, sp_id)
            # Change to hidden
            to_elem = vis_set.find(f".//{_qn('p:strVal')}")
            if to_elem is not None:
                to_elem.set("val", "hidden")
            child_tn.append(vis_set)

    elif effect == "pulse":
        # Pulse emphasis: use animScale (grow then shrink)
        scale = _build_anim_scale(ids, sp_id, duration, from_pct=100, to_pct=110)
        child_tn.append(scale)

    else:
        # Standard animEffect (fade, fly, wipe, float)
        filter_str = info.get("filter")
        if filter_str:
            # Replace direction placeholder
            if "{dir}" in filter_str:
                wipe_dir = WIPE_DIR.get(direction, "left")
                filter_str = filter_str.replace("{dir}", wipe_dir)

            transition = "in" if preset_class == "entr" else "out"
            anim_eff = _build_anim_effect(ids, sp_id, duration, transition, filter_str)
            child_tn.append(anim_eff)

    return par


# ---------------------------------------------------------------------------
# Timing Tree Builder
# ---------------------------------------------------------------------------

def build_timing_xml(animations, shape_map):
    """Build complete <p:timing> from animation spec list."""
    ids = IdCounter(1)

    # Root structure
    timing = etree.Element(_qn("p:timing"))
    tn_lst = etree.SubElement(timing, _qn("p:tnLst"))
    root_par = etree.SubElement(tn_lst, _qn("p:par"))

    root_ctn = etree.SubElement(root_par, _qn("p:cTn"), attrib={
        "id": ids.next(), "dur": "indefinite", "restart": "never", "nodeType": "tmRoot"
    })
    root_children = etree.SubElement(root_ctn, _qn("p:childTnLst"))

    # Main interactive sequence
    seq = etree.SubElement(root_children, _qn("p:seq"), attrib={
        "concurrent": "1", "nextAc": "seek"
    })
    seq_ctn = etree.SubElement(seq, _qn("p:cTn"), attrib={
        "id": ids.next(), "dur": "indefinite", "nodeType": "mainSeq"
    })
    seq_children = etree.SubElement(seq_ctn, _qn("p:childTnLst"))

    prev_cond = etree.SubElement(seq, _qn("p:prevCondLst"))
    etree.SubElement(prev_cond, _qn("p:cond"), attrib={"evt": "onPrev", "delay": "0"})
    next_cond = etree.SubElement(seq, _qn("p:nextCondLst"))
    etree.SubElement(next_cond, _qn("p:cond"), attrib={"evt": "onNext", "delay": "0"})

    # Track click points
    current_click_children = None
    grp_counter = 0

    for anim in animations:
        target = anim["target"]
        trigger = anim.get("trigger", "onClick")
        delay = str(anim.get("delay", 0))

        # Resolve target
        if isinstance(target, int):
            sp_id = target
        elif target in shape_map:
            sp_id = shape_map[target]
        else:
            print(f"WARNING: Shape '{target}' not found. Available: {list(shape_map.keys())}")
            continue

        # Build the effect par
        effect_par = build_effect_par(ids, sp_id, anim, grp_id=grp_counter)
        if effect_par is None:
            continue
        grp_counter += 1

        if trigger == "onClick":
            # New click point: outermost p:par > p:cTn > p:childTnLst
            click_par = etree.SubElement(seq_children, _qn("p:par"))
            click_ctn = etree.SubElement(click_par, _qn("p:cTn"), attrib={
                "id": ids.next(), "fill": "hold"
            })
            click_st = etree.SubElement(click_ctn, _qn("p:stCondLst"))
            etree.SubElement(click_st, _qn("p:cond"), attrib={"delay": "indefinite"})
            current_click_children = etree.SubElement(click_ctn, _qn("p:childTnLst"))

            # Set nodeType on the effect's cTn
            eff_ctn = effect_par.find(_qn("p:cTn"))
            if eff_ctn is not None:
                eff_ctn.set("nodeType", "clickEffect")

            # Wrap in delay par
            delay_par = etree.SubElement(current_click_children, _qn("p:par"))
            delay_ctn = etree.SubElement(delay_par, _qn("p:cTn"), attrib={
                "id": ids.next(), "fill": "hold"
            })
            delay_st = etree.SubElement(delay_ctn, _qn("p:stCondLst"))
            etree.SubElement(delay_st, _qn("p:cond"), attrib={"delay": delay})
            delay_children = etree.SubElement(delay_ctn, _qn("p:childTnLst"))
            delay_children.append(effect_par)

        elif trigger == "withPrevious":
            if current_click_children is None:
                click_par = etree.SubElement(seq_children, _qn("p:par"))
                click_ctn = etree.SubElement(click_par, _qn("p:cTn"), attrib={
                    "id": ids.next(), "fill": "hold"
                })
                click_st = etree.SubElement(click_ctn, _qn("p:stCondLst"))
                etree.SubElement(click_st, _qn("p:cond"), attrib={"delay": "0"})
                current_click_children = etree.SubElement(click_ctn, _qn("p:childTnLst"))

            eff_ctn = effect_par.find(_qn("p:cTn"))
            if eff_ctn is not None:
                eff_ctn.set("nodeType", "withEffect")

            delay_par = etree.SubElement(current_click_children, _qn("p:par"))
            delay_ctn = etree.SubElement(delay_par, _qn("p:cTn"), attrib={
                "id": ids.next(), "fill": "hold"
            })
            delay_st = etree.SubElement(delay_ctn, _qn("p:stCondLst"))
            etree.SubElement(delay_st, _qn("p:cond"), attrib={"delay": delay})
            delay_children = etree.SubElement(delay_ctn, _qn("p:childTnLst"))
            delay_children.append(effect_par)

        elif trigger == "afterPrevious":
            if current_click_children is None:
                click_par = etree.SubElement(seq_children, _qn("p:par"))
                click_ctn = etree.SubElement(click_par, _qn("p:cTn"), attrib={
                    "id": ids.next(), "fill": "hold"
                })
                click_st = etree.SubElement(click_ctn, _qn("p:stCondLst"))
                etree.SubElement(click_st, _qn("p:cond"), attrib={"delay": "0"})
                current_click_children = etree.SubElement(click_ctn, _qn("p:childTnLst"))

            eff_ctn = effect_par.find(_qn("p:cTn"))
            if eff_ctn is not None:
                eff_ctn.set("nodeType", "afterEffect")

            delay_par = etree.SubElement(current_click_children, _qn("p:par"))
            delay_ctn = etree.SubElement(delay_par, _qn("p:cTn"), attrib={
                "id": ids.next(), "fill": "hold"
            })
            delay_st = etree.SubElement(delay_ctn, _qn("p:stCondLst"))
            etree.SubElement(delay_st, _qn("p:cond"), attrib={"delay": delay})
            delay_children = etree.SubElement(delay_ctn, _qn("p:childTnLst"))
            delay_children.append(effect_par)

    return timing


# ---------------------------------------------------------------------------
# Shape Map + Operations
# ---------------------------------------------------------------------------

def get_shape_map(slide):
    smap = {}
    for shape in slide.shapes:
        smap[shape.name] = shape.shape_id
        smap[shape.shape_id] = shape.shape_id
    return smap


def list_shapes(pptx_path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        print(f"\n=== Slide {i + 1} ===")
        for shape in slide.shapes:
            shape_type = type(shape).__name__
            pos = f"({shape.left}, {shape.top})" if shape.left is not None else "(?, ?)"
            size = f"{shape.width}x{shape.height}" if shape.width is not None else "?x?"
            print(f"  id={shape.shape_id:3d}  name={shape.name!r:30s}  type={shape_type:20s}  pos={pos}  size={size}")


def inspect_animations(pptx_path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        slide_elem = slide._element
        timings = slide_elem.findall(f"{{{P_NS}}}timing")
        if timings:
            print(f"\n=== Slide {i + 1}: Has animations ===")
            for t in timings:
                xml_str = etree.tostring(t, pretty_print=True, encoding="unicode")
                print(xml_str[:3000])
        else:
            print(f"\n=== Slide {i + 1}: No animations ===")


def apply_animations(pptx_path, spec_path, output_path):
    prs = Presentation(pptx_path)
    with open(spec_path, "r") as f:
        spec = json.load(f)

    for slide_spec in spec.get("slides", []):
        slide_idx = slide_spec["slide_index"]
        if slide_idx >= len(prs.slides):
            print(f"WARNING: Slide index {slide_idx} out of range")
            continue

        slide = prs.slides[slide_idx]
        shape_map = get_shape_map(slide)
        animations = slide_spec.get("animations", [])

        if not animations:
            continue

        slide_elem = slide._element
        for existing in slide_elem.findall(f"{{{P_NS}}}timing"):
            slide_elem.remove(existing)

        timing_xml = build_timing_xml(animations, shape_map)
        slide_elem.append(timing_xml)
        print(f"Slide {slide_idx + 1}: Applied {len(animations)} animations")

        transition = slide_spec.get("transition")
        if transition:
            _apply_transition(slide, transition)

    prs.save(output_path)
    print(f"\nSaved: {output_path}")


def _apply_transition(slide, transition_spec):
    slide_elem = slide._element
    for existing in slide_elem.findall(f"{{{P_NS}}}transition"):
        slide_elem.remove(existing)

    trans_type = transition_spec.get("type", "fade")
    speed = transition_spec.get("speed", "med")

    trans_elem = etree.SubElement(slide_elem, _qn("p:transition"), attrib={"spd": speed})
    trans_types = {"fade": "fade", "push": "push", "wipe": "wipe", "split": "split", "cover": "cover"}
    if trans_type in trans_types:
        etree.SubElement(trans_elem, _qn(f"p:{trans_types[trans_type]}"))


def main():
    parser = argparse.ArgumentParser(description="PPTX Animation Engine")
    parser.add_argument("pptx", help="Input PPTX file")
    parser.add_argument("spec", nargs="?", help="Animation spec JSON file")
    parser.add_argument("-o", "--output", help="Output PPTX file")
    parser.add_argument("--list-shapes", action="store_true")
    parser.add_argument("--inspect", action="store_true")

    args = parser.parse_args()

    if args.list_shapes:
        list_shapes(args.pptx)
    elif args.inspect:
        inspect_animations(args.pptx)
    elif args.spec:
        output = args.output or args.pptx.replace(".pptx", "_animated.pptx")
        apply_animations(args.pptx, args.spec, output)
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
